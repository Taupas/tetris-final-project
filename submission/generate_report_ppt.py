# -*- coding: utf-8 -*-
"""Generate 7-9 slide final report PPT for Tetris (Flask x C)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
OUT_PPT = OUT_DIR / "期末報告_俄羅斯方塊.pptx"
OUT_DOWNLOADS = Path(r"C:\Users\USER\Downloads\期末報告_俄羅斯方塊.pptx")

TITLE_COLOR = RGBColor(0x1A, 0x3A, 0x6E)
BODY = RGBColor(0x33, 0x33, 0x33)
CODE_COLOR = RGBColor(0x1E, 0x1E, 0x1E)


def set_title(shape, text: str, size: int = 30) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR


def add_bullets(tf, lines: list[str], size: int = 17) -> None:
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        stripped = line.lstrip()
        if line.startswith("    "):
            p.level = 2
        elif line.startswith("  "):
            p.level = 1
        else:
            p.level = 0
        p.text = stripped
        p.font.size = Pt(size - p.level * 2)
        p.font.color.rgb = BODY


def slide_bullets(prs, title: str, lines: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s.shapes.title, title)
    add_bullets(s.placeholders[1].text_frame, lines)


def slide_code(prs, title: str, code: str, note: str = "") -> None:
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(s.shapes.title, title, 26)
    top = Inches(1.3)
    box = s.shapes.add_textbox(Inches(0.45), top, Inches(9.1), Inches(4.2))
    p = box.text_frame.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = Pt(13)
    p.font.color.rgb = CODE_COLOR
    if note:
        nb = s.shapes.add_textbox(Inches(0.45), Inches(5.6), Inches(9.1), Inches(1.2))
        add_bullets(nb.text_frame, [note], 15)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 1 封面 ──
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "俄羅斯方塊 Tetris"
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    s.placeholders[1].text = (
        "計算機程式語言 期末專題報告\n"
        "Flask 前端 × C 遊戲引擎（ctypes）\n"
        "http://127.0.0.1:5000/tetris"
    )

    # ── 2 架構 + Game Loop ──
    slide_bullets(prs, "1. 專案整體架構設計", [
        "【運作方式】網頁 Canvas 顯示 + C DLL 負責遊戲邏輯（非 SDL / Console）",
        "  瀏覽器 ←JSON→ Flask ←ctypes→ lib/tetris_engine.dll",
        "",
        "【Game Loop】",
        "  Input：鍵盤 ←↑↓→、空白鍵 → /tetris/api/input",
        "        拼圖模式：滑鼠拖放 → /tetris/api/block/place",
        "  Update：C 引擎 engine_tick — 下落、碰撞、消行、加分",
        "  Render：game.js / puzzle.js 依 API 狀態重繪 Canvas",
    ])

    # ── 3 模組化 ──
    slide_bullets(prs, "1. 模組化設計", [
        "【C 引擎 c_engine/】",
        "  board.c — 棋盤記憶體與格子讀寫",
        "  physics.c — 移動、旋轉、碰撞、鎖定",
        "  tetromino.c — 七種方塊 I/O/T/L/J/S/Z",
        "  queue.c — 7-bag 出塊佇列（加分）",
        "  engine_api.c — 經典模式 10×20",
        "  block_puzzle_api.c — 拼圖模式 8×8",
        "【其他】engine_bridge.py（ctypes）· app.py（路由）· game.js（繪圖）",
    ])

    # ── 4 Pointer & Malloc ──
    slide_code(
        prs,
        "2. Pointer 與 Malloc — 動態遊戲地圖",
        (
            "Board *board_create(int rows, int cols) {\n"
            "    Board *b = (Board *)malloc(sizeof(Board));\n"
            "    b->cells = (int *)malloc(rows * cols * sizeof(int));\n"
            "    // 失敗時 free(b)，避免 leak\n"
            "    return b;\n"
            "}\n"
            "// 一維 cells[row*cols+col] 模擬 20x10 二維棋盤\n"
            "int board_get(const Board *b, int row, int col);"
        ),
        "★ Game Over / engine_destroy 時呼叫 board_destroy → free(cells) 再 free(b)",
    )

    # ── 5 Pass by Reference ──
    slide_code(
        prs,
        "2. 指標傳遞遊戲狀態（Pass by Reference）",
        (
            "int physics_move(Board *b, Piece *p, int dx, int dy) {\n"
            "    Piece trial = *p;\n"
            "    trial.x += dx;  trial.y += dy;\n"
            "    if (physics_can_place(b, &trial)) {\n"
            "        *p = trial;   // 透過指標直接修改方塊\n"
            "        return 1;\n"
            "    }\n"
            "    return 0;\n"
            "}\n"
            "void engine_tick(int game_id);  // Game* g = game_slots[id]"
        ),
        "避免複製整個 Game struct；用 -> 運算子修改堆積上的遊戲狀態",
    )

    # ── 6 資料結構 ──
    slide_bullets(prs, "3. 資料結構的應用與設計理由", [
        "【二維陣列概念 Grid】",
        "  cells[row*cols+col] 表示 10×20 / 8×8 盤面",
        "  SHAPES[type][rot] 表示 7 種 Tetromino 旋轉形狀",
        "  快速判定：空格？整行滿？可否放置？",
        "",
        "【結構體 Struct】",
        "  Piece { type, rotation, x, y } — 當前方塊",
        "  Game { Board*, Piece, IntQueue*, score, level } — 整局狀態",
        "",
        "【環形 Queue】IntQueue — 7-bag 公平隨機出塊（加分項目）",
    ])

    # ── 7 Cursor 協作 ──
    slide_bullets(prs, "4. 與 Cursor 協作解決的具體問題", [
        "定位：Cursor 是高效助手，我是主導的架構師",
        "",
        "① 旋轉與踢牆（Wall Kick）",
        "  physics_rotate + try_rotate 檢查旋轉後是否撞牆",
        "",
        "② 記憶體 / ctypes 除錯",
        "  拼圖方塊不顯示 → 修正 slot_dr 扁平陣列讀取錯位",
        "",
        "③ 程式碼重構",
        "  拆成 board / physics / engine_api 等模組，Flask 僅路由",
        "",
        "④ 遊戲邏輯：game over 判定、L 形貼邊 resolve_placement",
    ])

    # ── 8 總結 ──
    slide_bullets(prs, "總結", [
        "三層架構：Browser → Flask → C DLL（ctypes）",
        "Malloc/Free 對稱配置，Game Over 時完整釋放記憶體",
        "Grid + Struct + Queue 讓遊戲邏輯清晰可維護",
        "Cursor 加速實作與除錯，架構決策由開發者主導",
        "",
        "Demo：http://127.0.0.1:5000/tetris",
        "GitHub：submission/ 內含 README 與完整文件",
    ])

    # ── 9 Q&A ──
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "謝謝聆聽"
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    s.placeholders[1].text = "Q & A"

    return prs


if __name__ == "__main__":
    ppt = build()
    ppt.save(str(OUT_PPT))
    ppt.save(str(OUT_DOWNLOADS))
    print(f"共 {len(ppt.slides)} 頁")
    print(f"已產生：{OUT_PPT}")
    print(f"已複製：{OUT_DOWNLOADS}")
